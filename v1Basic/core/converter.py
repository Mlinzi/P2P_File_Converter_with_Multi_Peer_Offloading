"""
converter.py — handles all file conversion logic
Supports: documents (LibreOffice/docx2pdf), images (Pillow), audio/video (FFmpeg), PDF combining (pypdf)
"""

import subprocess
import shutil
from pathlib import Path

# ---------------------------------------------------------------------------
# Tool discovery
# ---------------------------------------------------------------------------

def _find_ffmpeg():
    bundled = Path(__file__).parent.parent / 'bin' / 'ffmpeg.exe'
    if bundled.exists():
        return str(bundled)
    if shutil.which('ffmpeg'):
        return 'ffmpeg'
    try:
        import imageio_ffmpeg
        return imageio_ffmpeg.get_ffmpeg_exe()
    except Exception:
        return None


def _find_soffice():
    """Find LibreOffice soffice binary."""
    # Common Windows install paths
    candidates = [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
    ]
    for path in candidates:
        if Path(path).exists():
            return path
    # Try PATH (Linux/Mac or if added to PATH on Windows)
    if shutil.which('soffice'):
        return 'soffice'
    return None


FFMPEG  = _find_ffmpeg()
SOFFICE = _find_soffice()

def _has_docx2pdf():
    try: import docx2pdf; return True
    except ImportError: return False

HAS_DOCX2PDF = _has_docx2pdf()

def _has_comtypes() -> bool:
    try: import comtypes.client; return True
    except ImportError: return False

HAS_COMTYPES = _has_comtypes()   # PowerPoint / Excel COM (Windows only)

def _detect_gpu_encoders():
    """Returns available HW video encoders in priority order: NVENC > AMF > QSV."""
    if not FFMPEG: return []
    try:
        r = subprocess.run([FFMPEG, '-hide_banner', '-encoders'],
                           capture_output=True, text=True, timeout=5)
        return [e for e in ['h264_nvenc', 'h264_amf', 'h264_qsv'] if e in r.stdout]
    except Exception:
        return []

GPU_ENCODERS = _detect_gpu_encoders()
USE_GPU      = False   # toggled at runtime

def set_gpu_accel(enabled: bool):
    global USE_GPU
    USE_GPU = enabled

def get_gpu_accel() -> bool:
    return USE_GPU

# ---------------------------------------------------------------------------
# Format definitions — grouped by document type
# ---------------------------------------------------------------------------

# Word processor formats
WORD_FORMATS         = {'docx', 'doc', 'odt', 'rtf'}
# Presentation formats
PRESENTATION_FORMATS = {'pptx', 'ppt', 'odp'}
# Spreadsheet formats
SPREADSHEET_FORMATS  = {'xlsx', 'xls', 'ods', 'csv'}

ALL_DOCUMENT_FORMATS = WORD_FORMATS | PRESENTATION_FORMATS | SPREADSHEET_FORMATS

IMAGE_FORMATS = {'png', 'jpeg', 'bmp', 'gif', 'tiff', 'webp'}
AUDIO_FORMATS = {'mp3', 'wav', 'flac', 'ogg', 'm4a', 'aac'}
VIDEO_FORMATS = {'mp4', 'avi', 'mkv', 'mov', 'webm'}

# LibreOffice filter names for --convert-to (needed for MS formats)
LIBREOFFICE_FILTERS = {
    'docx': 'docx:"MS Word 2007 XML"',
    'xlsx': 'xlsx:"Calc MS Excel 2007 XML"',
    'pptx': 'pptx:"Impress MS PowerPoint 2007 XML"',
    'csv' : 'csv:"Text - txt - csv (StarCalc)"',
    'txt' : 'txt:Text',
}

# ---------------------------------------------------------------------------
# Conversion map — what each format can be converted to
# ---------------------------------------------------------------------------

def _doc_targets(fmt, group):
    """All formats in same group (minus self) + pdf."""
    return sorted((group - {fmt}) | {'pdf'})

CONVERSION_MAP = {
    # Word processor ↔ each other + PDF
    **{fmt: _doc_targets(fmt, WORD_FORMATS) for fmt in WORD_FORMATS},
    # Presentations ↔ each other + PDF
    **{fmt: _doc_targets(fmt, PRESENTATION_FORMATS) for fmt in PRESENTATION_FORMATS},
    # Spreadsheets ↔ each other + PDF
    **{fmt: _doc_targets(fmt, SPREADSHEET_FORMATS) for fmt in SPREADSHEET_FORMATS},
    # PDF: only combining (handled separately, not via convert())
    'pdf': [],
    # Images ↔ any other image
    **{fmt: sorted(IMAGE_FORMATS - {fmt}) for fmt in IMAGE_FORMATS},
    # Audio ↔ any other audio
    **{fmt: sorted(AUDIO_FORMATS - {fmt}) for fmt in AUDIO_FORMATS},
    # Video ↔ any other video + mp3 (audio extract)
    **{fmt: sorted(VIDEO_FORMATS - {fmt}) + ['mp3'] for fmt in VIDEO_FORMATS},
}

# ---------------------------------------------------------------------------
# Public helpers
# ---------------------------------------------------------------------------

def detect_format(filepath: str) -> str:
    """Return normalised lowercase extension. Treats .jpg as jpeg."""
    ext = Path(filepath).suffix.lstrip('.').lower()
    if ext == 'jpg':
        return 'jpeg'
    return ext


def get_available_outputs(input_format: str) -> list[str]:
    """Return only the formats actually achievable with currently installed tools."""
    fmt     = input_format.lower()
    targets = CONVERSION_MAP.get(fmt, [])
    result  = []
    for t in targets:
        if fmt in AUDIO_FORMATS or fmt in VIDEO_FORMATS:
            if FFMPEG: result.append(t)           # audio/video needs FFmpeg
        elif fmt in ALL_DOCUMENT_FORMATS:
            if t == 'pdf':
                can = SOFFICE
                can = can or (HAS_DOCX2PDF  and fmt in WORD_FORMATS)          # docx2pdf (Word)
                can = can or (HAS_COMTYPES  and fmt in {'pptx', 'ppt'})       # PowerPoint COM
                can = can or fmt in {'pptx', 'ppt'}                           # python-pptx fallback
                if can: result.append(t)
            else:
                if SOFFICE: result.append(t)      # format-to-format needs LibreOffice
        else:
            result.append(t)                      # images only need Pillow
    return result


# ---------------------------------------------------------------------------
# Main conversion entry point
# ---------------------------------------------------------------------------

def convert(input_path: str, output_format: str, output_dir: str,
            use_gpu: bool | None = None) -> Path:
    """
    Convert a single file to output_format.
    use_gpu: override GPU flag for this call only (None = use global USE_GPU setting).
    Returns path to the converted output file.
    Raises ValueError for unsupported combos, RuntimeError on tool failure.
    """
    input_path    = Path(input_path)
    output_dir    = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    input_fmt     = detect_format(str(input_path))
    output_format = output_format.lower()
    if output_format == 'jpg':
        output_format = 'jpeg'

    valid = get_available_outputs(input_fmt)
    if output_format not in valid:
        raise ValueError(
            f"Cannot convert {input_fmt.upper()} → {output_format.upper()}. "
            f"Valid targets: {valid}"
        )

    disk_ext    = 'jpg' if output_format == 'jpeg' else output_format
    output_path = output_dir / f"{input_path.stem}_converted.{disk_ext}"

    effective_gpu = USE_GPU if use_gpu is None else use_gpu

    if input_fmt in ALL_DOCUMENT_FORMATS:
        return _convert_document(input_path, output_format, output_dir, output_path)

    if input_fmt in IMAGE_FORMATS:
        return _convert_image(input_path, output_path, output_format)

    if input_fmt in AUDIO_FORMATS or input_fmt in VIDEO_FORMATS:
        return _convert_ffmpeg(input_path, output_path, effective_gpu)

    raise ValueError(f"Unsupported input format: {input_fmt}")


def combine_pdfs(input_paths: list[str], output_path: str) -> Path:
    """Merge multiple PDFs into one."""
    from pypdf import PdfMerger

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    merger = PdfMerger()
    for p in input_paths:
        merger.append(str(p))
    merger.write(str(output_path))
    merger.close()
    return output_path


# ---------------------------------------------------------------------------
# Internal converters
# ---------------------------------------------------------------------------

def _convert_document(input_path: Path, output_format: str, output_dir: Path, output_path: Path) -> Path:
    """
    PDF output priority:
      Word formats  → docx2pdf (Word COM)  → LibreOffice
      pptx/ppt      → LibreOffice           → PowerPoint COM
      everything else → LibreOffice
    Other formats: LibreOffice only.
    """
    input_fmt = detect_format(str(input_path))

    if output_format == 'pdf':
        # Word formats: try docx2pdf (Word COM) first
        if input_fmt in WORD_FORMATS:
            try:
                from docx2pdf import convert as _docx2pdf
                _docx2pdf(str(input_path), str(output_dir))
                produced = output_dir / f"{input_path.stem}.pdf"
                if produced.exists():
                    produced.rename(output_path)
                    return output_path
            except Exception:
                pass

        # pptx/ppt: try LibreOffice → PowerPoint COM → python-pptx fallback
        if input_fmt in {'pptx', 'ppt'}:
            if SOFFICE:
                try:
                    return _convert_with_soffice(input_path, output_format, output_dir, output_path)
                except Exception:
                    pass
            if HAS_COMTYPES:
                return _convert_pptx_com(input_path, output_path)
            try:
                return _convert_pptx_python(input_path, output_path)
            except Exception as e:
                raise RuntimeError(
                    f"Cannot convert PowerPoint to PDF: install LibreOffice, Microsoft PowerPoint, "
                    f"or run: pip install python-pptx pillow reportlab\n({e})"
                )

    # LibreOffice for everything else
    return _convert_with_soffice(input_path, output_format, output_dir, output_path)


def _convert_pptx_com(input_path: Path, output_path: Path) -> Path:
    """Convert pptx/ppt → PDF using PowerPoint COM (Windows, requires MS PowerPoint)."""
    import comtypes
    import comtypes.client

    ppt_path   = str(input_path.resolve())
    powerpoint = None
    presentation = None
    try:
        comtypes.CoInitialize()
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.DisplayAlerts = 0
        presentation = powerpoint.Presentations.Open(
            ppt_path, ReadOnly=True, Untitled=False, WithWindow=False
        )
        # PowerPoint adds .pdf itself — pass path without extension
        pdf_base = str(output_path.with_suffix(''))
        presentation.SaveAs(pdf_base, 32)   # 32 = ppSaveAsPDF
        produced = output_path.with_suffix('.pdf')
        if produced != output_path and produced.exists():
            produced.rename(output_path)
        return output_path
    except Exception as e:
        raise RuntimeError(f"PowerPoint COM failed: {e}")
    finally:
        try:
            if presentation: presentation.Close()
        except Exception: pass
        try:
            if powerpoint: powerpoint.Quit()
        except Exception: pass
        try:
            comtypes.CoUninitialize()
        except Exception: pass


def _convert_pptx_python(input_path: Path, output_path: Path) -> Path:
    """
    Pure-Python pptx → PDF fallback using python-pptx + Pillow + reportlab.
    Renders each slide as an image and combines into a PDF.
    Quality is basic (no fonts/transitions) but works with no external tools.
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw
    from reportlab.pdfgen import canvas
    import io

    prs = Presentation(str(input_path))

    slide_w = prs.slide_width  / 914400
    slide_h = prs.slide_height / 914400

    dpi      = 150
    px_w     = int(slide_w * dpi)
    px_h     = int(slide_h * dpi)
    page_size = (slide_w * 72, slide_h * 72)

    c = canvas.Canvas(str(output_path), pagesize=page_size)

    for slide in prs.slides:
        img  = Image.new('RGB', (px_w, px_h), 'white')
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            if shape.has_text_frame:
                x = int(shape.left / 914400 * dpi)
                y = int(shape.top  / 914400 * dpi)
                for para in shape.text_frame.paragraphs:
                    txt = para.text
                    if txt.strip():
                        draw.text((x, y), txt, fill='black')
                        y += 18

        buf = io.BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)

        img_tmp = output_path.with_suffix('.tmp_slide.png')
        Image.open(buf).save(str(img_tmp))

        c.drawImage(str(img_tmp), 0, 0, width=page_size[0], height=page_size[1])
        c.showPage()
        try: img_tmp.unlink()
        except Exception: pass

    c.save()
    return output_path


def _convert_with_soffice(input_path: Path, output_format: str, output_dir: Path, output_path: Path) -> Path:
    if not SOFFICE:
        raise RuntimeError(
            "LibreOffice not found. Install LibreOffice (https://www.libreoffice.org/download) "
            "or add soffice to PATH."
        )

    filter_arg = LIBREOFFICE_FILTERS.get(output_format, output_format)

    cmd = [
        SOFFICE,
        '--headless',
        '--convert-to', filter_arg,
        '--outdir', str(output_dir),
        str(input_path),
    ]

    result = subprocess.run(cmd, capture_output=True, text=True)

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed:\n{result.stderr[-1000:]}")

    # LibreOffice writes <stem>.<ext> in output_dir
    disk_ext  = 'jpg' if output_format == 'jpeg' else output_format
    produced  = output_dir / f"{input_path.stem}.{disk_ext}"

    if not produced.exists():
        raise RuntimeError(f"LibreOffice ran but output not found: {produced}")

    produced.rename(output_path)
    return output_path


def _convert_image(input_path: Path, output_path: Path, output_format: str) -> Path:
    try:
        from PIL import Image
    except ImportError:
        raise RuntimeError("Pillow not installed. Run: pip install pillow")

    img = Image.open(input_path)

    if img.mode == 'RGBA' and output_format == 'jpeg':
        img = img.convert('RGB')

    if getattr(img, 'is_animated', False) and output_format != 'gif':
        img.seek(0)

    pil_fmt = 'JPEG' if output_format == 'jpeg' else output_format.upper()
    img.save(output_path, format=pil_fmt)
    return output_path


def _convert_ffmpeg(input_path: Path, output_path: Path,
                    use_gpu: bool = False) -> Path:
    if not FFMPEG:
        raise RuntimeError(
            "FFmpeg not found. Place ffmpeg.exe in /bin/ or install FFmpeg system-wide."
        )

    out_fmt = output_path.suffix.lstrip('.').lower()
    cmd     = [FFMPEG, '-i', str(input_path)]

    # GPU encoding: only if requested, a HW encoder exists, and output is video
    if use_gpu and GPU_ENCODERS and out_fmt in VIDEO_FORMATS:
        cmd += ['-c:v', GPU_ENCODERS[0]]   # e.g. h264_nvenc / h264_amf / h264_qsv

    cmd += ['-y', str(output_path)]

    result = subprocess.run(cmd, capture_output=True, text=True)

    if result.returncode != 0:
        # GPU encode failed — retry with software encoder
        if use_gpu and GPU_ENCODERS and out_fmt in VIDEO_FORMATS:
            result = subprocess.run(
                [FFMPEG, '-i', str(input_path), '-y', str(output_path)],
                capture_output=True, text=True
            )
        if result.returncode != 0:
            raise RuntimeError(f"FFmpeg failed:\n{result.stderr[-1000:]}")

    return output_path


# ---------------------------------------------------------------------------
# Quick standalone test
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    print("CONVERSION_MAP:")
    for fmt, targets in sorted(CONVERSION_MAP.items()):
        print(f"  {fmt:<8} → {targets}")

    ffmpeg_status  = f"found ({FFMPEG})"  if FFMPEG  else "NOT FOUND"
    soffice_status = f"found ({SOFFICE})" if SOFFICE else "NOT FOUND"
    print(f"\nTool status:  ffmpeg={ffmpeg_status}  soffice={soffice_status}")
    print(f"\ndetect_format('report.DOCX') = '{detect_format('report.DOCX')}'")
    print(f"docx  → {get_available_outputs('docx')}")
    print(f"pptx  → {get_available_outputs('pptx')}")
    print(f"xlsx  → {get_available_outputs('xlsx')}")
    print(f"mp4   → {get_available_outputs('mp4')}")
    print(f"pdf   → {get_available_outputs('pdf')} (combine only)")
